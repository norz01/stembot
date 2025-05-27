import streamlit as st
import requests
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from fpdf import FPDF # Pastikan ini adalah fpdf2: pip install fpdf2
import pandas as pd
from datetime import datetime
import os
import json
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
import time
from PIL import Image
import pytesseract # Anda mungkin perlu memasang Tesseract OCR: https://github.com/tesseract-ocr/tesseract
import fitz  # PyMuPDF untuk PDF: pip install PyMuPDF

# --- KONFIGURASI ---
HISTORY_DIR = "chat_sessions"
UPLOAD_DIR = "uploaded_files" # Direktori untuk fail yang dimuat naik
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
DEFAULT_OLLAMA_MODEL = os.getenv("DEFAULT_OLLAMA_MODEL", "llama3") # Model lalai

# Konfigurasi untuk ciri dari chatbot2
LOGO_PATH = os.getenv("ikm_logo", "ikm_logo.png") # Letakkan logo anda di sini dan namakannya ikm_logo.png atau set pembolehubah persekitaran
WATERMARK_TEXT = os.getenv("CHATBOT_WATERMARK_TEXT", "IKM Besut")
# Pastikan Tesseract OCR dipasang dan dikonfigurasi dalam PATH sistem anda, atau setkan laluan tesseract_cmd
# Contoh: pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

os.makedirs(HISTORY_DIR, exist_ok=True)
os.makedirs(UPLOAD_DIR, exist_ok=True)

# --- FUNGSI HELPER (Gabungan dan Penambahbaikan) ---

@st.cache_data(ttl=300)
def get_ollama_models_cached():
    """Mendapatkan senarai model yang tersedia dari Ollama dan mengcache hasilnya."""
    try:
        response = requests.get(f'{OLLAMA_BASE_URL}/api/tags', timeout=10)
        response.raise_for_status()
        models_data = response.json().get('models', [])
        if not models_data:
            return []
        return sorted([model['name'] for model in models_data])
    except requests.exceptions.Timeout:
        st.error(f"Gagal mendapatkan senarai model: Permintaan ke Ollama tamat masa.")
        return []
    except requests.exceptions.RequestException:
        return []
    except KeyError:
        st.error("Format respons senarai model tidak dijangka dari Ollama.")
        return []

# Namakan semula fungsi asal
def query_ollama_non_stream(prompt, chat_history, selected_model):
    """Menghantar pertanyaan ke Ollama dan mengembalikan respons serta masa penjanaan (NON-STREAM)."""
    messages_for_api = [{"role": msg["role"], "content": msg["content"]} for msg in chat_history]
    if not messages_for_api or messages_for_api[-1]["content"] != prompt or messages_for_api[-1]["role"] != "user":
         messages_for_api.append({"role": "user", "content": prompt})

    start_time = time.time()
    try:
        payload = {'model': selected_model, 'messages': messages_for_api, 'stream': False} # STREAM FALSE
        # Anda mungkin mahu menggunakan OLLAMA_REQUEST_TIMEOUT di sini juga jika anda melaksanakannya
        response = requests.post(f'{OLLAMA_BASE_URL}/api/chat', json=payload, timeout=600) # Timeout asal
        response.raise_for_status()
        
        end_time = time.time()
        processing_time = end_time - start_time
        full_response_data = response.json()
        assistant_reply = full_response_data.get('message', {}).get('content', "Maaf, saya tidak dapat respons yang betul.")
        return assistant_reply, processing_time

    except requests.exceptions.Timeout:
        end_time = time.time(); processing_time = end_time - start_time
        st.error(f"Gagal mendapatkan respons: Permintaan ke Ollama tamat masa selepas {processing_time:.2f}s.")
        return "Maaf, permintaan tamat masa.", processing_time
    except requests.exceptions.RequestException as e:
        end_time = time.time(); processing_time = end_time - start_time
        st.error(f"Masalah menyambung ke Ollama: {e} (selepas {processing_time:.2f}s)")
        return "Maaf, berlaku masalah semasa menghubungi Ollama.", processing_time
    except KeyError:
        end_time = time.time(); processing_time = end_time - start_time
        st.error(f"Format respons tidak dijangka dari Ollama (selepas {processing_time:.2f}s).")
        return "Maaf, format respons dari Ollama tidak seperti yang dijangkakan.", processing_time

# Fungsi baru untuk strim
def query_ollama(prompt, chat_history, selected_model, response_placeholder): # Tambah response_placeholder
    """Menghantar pertanyaan ke Ollama dan stream respons ke placeholder Streamlit."""
    messages_for_api = [{"role": msg["role"], "content": msg["content"]} for msg in chat_history]
    # Jika prompt bukan sebahagian daripada mesej terakhir dalam chat_history, tambahkannya
    # Ini penting kerana chat_history yang dihantar mungkin sudah termasuk prompt pengguna terkini
    # Kita hanya mahu memastikan ia ada untuk API.
    # Untuk strim, prompt pengguna sudah ada dalam chat_history yang dihantar dari main().

    start_time = time.time()
    full_response_content = ""
    try:
        payload = {'model': selected_model, 'messages': messages_for_api, 'stream': True}
        # Anda mungkin mahu menggunakan OLLAMA_REQUEST_TIMEOUT di sini juga jika anda melaksanakannya
        with requests.post(f'{OLLAMA_BASE_URL}/api/chat', json=payload, stream=True, timeout=600) as response: # Timeout asal
            response.raise_for_status()
            for line in response.iter_lines():
                if line:
                    decoded_line = line.decode('utf-8')
                    try:
                        chunk = json.loads(decoded_line)
                        # Semak jika 'message' wujud dan ia adalah kamus
                        if 'message' in chunk and isinstance(chunk['message'], dict):
                            content_piece = chunk['message'].get('content', '')
                            if content_piece:
                                full_response_content += content_piece
                                response_placeholder.markdown(full_response_content + "▌")
                        
                        if chunk.get("done"): # Semak status 'done'
                            # Kadang-kadang mesej terakhir ada dalam chunk 'done' itu sendiri
                            # atau ia menandakan akhir strim
                            final_message_in_done_chunk = chunk.get('message', {}).get('content', '')
                            if final_message_in_done_chunk and not full_response_content.endswith(final_message_in_done_chunk):
                                full_response_content += final_message_in_done_chunk
                            break # Keluar dari gelung iter_lines apabila 'done' adalah True
                    except json.JSONDecodeError:
                        # Abaikan baris yang tidak dapat diproses sebagai JSON, mungkin baris kosong atau metadata lain
                        # st.warning(f"Gagal memproses baris strim: {decoded_line}")
                        pass
        
        end_time = time.time()
        processing_time = end_time - start_time
        response_placeholder.markdown(full_response_content) # Papar respons akhir tanpa kursor
        return full_response_content, processing_time

    except requests.exceptions.Timeout:
        end_time = time.time(); processing_time = end_time - start_time
        error_message = f"Gagal mendapatkan respons: Permintaan ke Ollama tamat masa selepas {processing_time:.2f}s."
        response_placeholder.error(error_message)
        return "Maaf, permintaan tamat masa.", processing_time
    except requests.exceptions.RequestException as e:
        end_time = time.time(); processing_time = end_time - start_time
        error_message = f"Masalah menyambung ke Ollama: {e} (selepas {processing_time:.2f}s)"
        response_placeholder.error(error_message)
        return "Maaf, berlaku masalah semasa menghubungi Ollama.", processing_time
    except Exception as e: # Tangkap ralat lain yang mungkin berlaku
        end_time = time.time(); processing_time = end_time - start_time
        error_message = f"Ralat tidak dijangka semasa strim: {e} (selepas {processing_time:.2f}s)"
        response_placeholder.error(error_message)
        return "Maaf, ralat tidak dijangka berlaku.", processing_time

def save_chat_session(session_id, history):
    filepath = os.path.join(HISTORY_DIR, f"{session_id}.json")
    try:
        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(history, f, indent=2)
    except IOError as e:
        st.error(f"Gagal menyimpan sesi Perbualan '{session_id}': {e}")

def load_chat_session(session_id):
    filepath = os.path.join(HISTORY_DIR, f"{session_id}.json")
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return json.load(f)
    except FileNotFoundError:
        return []
    except (json.JSONDecodeError, IOError) as e:
        st.error(f"Gagal memuatkan atau membaca sesi Perbualan '{session_id}': {e}")
        return []

def load_all_session_ids():
    try:
        files = [f.replace(".json", "") for f in os.listdir(HISTORY_DIR) if f.endswith(".json")]
        def sort_key(filename):
            try:
                parts = filename.split('_')
                if len(parts) >= 2:
                    return datetime.strptime(f"{parts[0]}_{parts[1]}", "%Y%m%d_%H%M%S")
            except (ValueError, IndexError): pass
            return datetime.min # Fallback untuk nama fail yang tidak mengikut format tarikh
        return sorted(files, key=sort_key, reverse=True)
    except OSError as e:
        st.error(f"Gagal membaca direktori sesi: {e}")
        return []

def delete_chat_session_file(session_id):
    filepath = os.path.join(HISTORY_DIR, f"{session_id}.json")
    try:
        if os.path.exists(filepath):
            os.remove(filepath)
            st.success(f"Sesi Perbualan '{session_id}' berjaya dipadam.")
            return True
        else:
            st.warning(f"Fail sesi Perbualan '{session_id}' tidak ditemui untuk dipadam.")
            return False
    except OSError as e:
        st.error(f"Gagal memadam sesi Perbualan '{session_id}': {e}")
        return False

def delete_all_chat_sessions():
    deleted_count = 0; errors = []
    try:
        for filename in os.listdir(HISTORY_DIR):
            if filename.endswith(".json"):
                filepath = os.path.join(HISTORY_DIR, filename)
                try: os.remove(filepath); deleted_count += 1
                except OSError as e: errors.append(f"Gagal memadam {filename}: {e}")
        if errors:
            for error in errors: st.error(error)
        if deleted_count > 0: st.success(f"{deleted_count} sesi Perbualan berjaya dipadam.")
        else: st.info("Tiada sesi Perbualan ditemui untuk dipadam.")
        return True
    except OSError as e:
        st.error(f"Gagal mengakses direktori sesi: {e}")
        return False

# --- FUNGSI EKSTRAKSI TEKS DARI FAIL (dari chatbot2) ---
def extract_text_from_file(uploaded_file_obj):
    extracted_text = ""
    filename = uploaded_file_obj.name
    file_bytes = uploaded_file_obj.getvalue() # Dapatkan bytes terus
    temp_docx_path = None # Inisialisasi di luar untuk rujukan dalam except utama jika perlu

    try:
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
            image = Image.open(uploaded_file_obj) # PIL boleh buka objek fail secara terus
            extracted_text = pytesseract.image_to_string(image)
            if not extracted_text.strip():
                 st.info(f"Tiada teks dapat diekstrak dari imej '{filename}' menggunakan OCR.")
        elif filename.lower().endswith(".txt"):
            extracted_text = file_bytes.decode('utf-8', errors='ignore')
        elif filename.lower().endswith(".docx"):
            # Simpan sementara untuk dibaca oleh python-docx
            temp_docx_path = os.path.join(UPLOAD_DIR, f"temp_{filename}") # Tetapkan temp_docx_path di sini
            try:
                with open(temp_docx_path, "wb") as f:
                    f.write(file_bytes)
                doc = Document(temp_docx_path)
                extracted_text = "\n".join([para.text for para in doc.paragraphs])
            except Exception as e_docx:
                st.error(f"Ralat semasa memproses fail DOCX '{filename}': {e_docx}")
                # Biarkan extracted_text sebagai string kosong atau set kepada None jika lebih sesuai
                # extracted_text = None # Jika anda mahu return None secara eksplisit pada ralat DOCX
            finally:
                if temp_docx_path and os.path.exists(temp_docx_path): # Periksa jika temp_docx_path telah ditetapkan
                    try:
                        os.remove(temp_docx_path)
                    except OSError as e_remove:
                        st.warning(f"Gagal memadam fail sementara DOCX '{temp_docx_path}': {e_remove}")
        elif filename.lower().endswith(".pdf"):
            # PyMuPDF boleh buka bytes terus
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            for page in doc:
                extracted_text += page.get_text()
            doc.close()
        else:
            st.warning(f"Jenis fail '{filename}' tidak disokong untuk ekstraksi teks.")
            return None # Return None jika jenis fail tidak disokong

        # Hanya strip jika extracted_text adalah string
        return extracted_text.strip() if isinstance(extracted_text, str) else None

    except Exception as e:
        st.error(f"Ralat umum semasa memproses fail '{filename}': {e}")
        # Cuba padam fail sementara DOCX jika ia dicipta dan ralat berlaku di luar blok DOCX
        if temp_docx_path and os.path.exists(temp_docx_path): # Periksa jika temp_docx_path telah ditetapkan
            try:
                os.remove(temp_docx_path)
            except OSError as e_remove_outer:
                st.warning(f"Gagal memadam fail sementara DOCX (luar) '{temp_docx_path}': {e_remove_outer}")
        return None

# --- FUNGSI EKSPORT (Gabungan dengan logo/watermark dari chatbot2) ---
def format_conversation_text(chat_history, include_user=True, include_assistant=True):
    lines = []
    for msg in chat_history:
        if (msg["role"] == "user" and include_user) or \
           (msg["role"] == "assistant" and include_assistant):
            lines.append(f"{msg['role'].capitalize()}: {msg['content'].strip()}")
    return "\n\n".join(lines)

def save_to_word(text_content, filename='output.docx', logo_path=LOGO_PATH, watermark_text=WATERMARK_TEXT):
    doc = Document()
    if logo_path and os.path.exists(logo_path):
        try:
            paragraph = doc.add_paragraph()
            run = paragraph.add_run()
            run.add_picture(logo_path, width=DocxInches(2.0)) # Saiz logo boleh laras
            paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            doc.add_paragraph() # Baris kosong selepas logo
        except Exception as e:
            st.warning(f"Gagal menambah logo pada Word: {e}. Pastikan fail imej sah.")

    if watermark_text:
        watermark_para = doc.add_paragraph()
        run = watermark_para.add_run(watermark_text)
        font = run.font
        font.size = DocxPt(36) # Saiz tera air
        font.color.rgb = DocxRGBColor(192, 192, 192)  # Kelabu cair
        font.bold = True
        watermark_para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        doc.add_paragraph()

    for para_block in text_content.split("\n\n"):
        doc.add_paragraph(para_block.strip())
    try: doc.save(filename); return True
    except IOError as e: st.error(f"Gagal menyimpan ke Word: {e}"); return False

def save_to_pdf(text_content, filename='output.pdf', logo_path=LOGO_PATH, watermark_text=WATERMARK_TEXT):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    # --- PERUBAHAN BERMULA DI SINI UNTUK SOKONGAN UNICODE ---
    FONT_DIR = "fonts"  # Pastikan direktori ini wujud
    # Nama fail fon. Anda boleh menamakannya apa sahaja, tetapi pastikan ia sepadan dengan fail sebenar.
    FONT_REGULAR_FILENAME = "DejaVuSans.ttf"
    FONT_REGULAR_PATH = os.path.join(FONT_DIR, FONT_REGULAR_FILENAME)

    # Nama keluarga fon yang akan digunakan dalam FPDF
    UNICODE_FONT_FAMILY = "DejaVuSans"
    DEFAULT_FALLBACK_FONT = "Arial" # Fon lalai FPDF jika fon Unicode gagal dimuatkan

    current_font_family_for_content = DEFAULT_FALLBACK_FONT
    current_font_family_for_watermark = DEFAULT_FALLBACK_FONT
    watermark_style = 'B' # Gaya Bold untuk watermark jika guna fon lalai

    if os.path.exists(FONT_REGULAR_PATH):
        try:
            # Tambah fon Unicode. '' untuk gaya biasa. uni=True adalah penting.
            pdf.add_font(UNICODE_FONT_FAMILY, '', FONT_REGULAR_PATH, uni=True)
            current_font_family_for_content = UNICODE_FONT_FAMILY
            current_font_family_for_watermark = UNICODE_FONT_FAMILY
            watermark_style = '' # Fon Unicode mungkin sudah cukup tebal atau tidak memerlukan gaya 'B' eksplisit
                                 # Jika anda mempunyai versi Bold (cth., DejaVuSans-Bold.ttf), anda boleh menambahkannya secara berasingan:
                                 # pdf.add_font(UNICODE_FONT_FAMILY, 'B', "path/to/DejaVuSans-Bold.ttf", uni=True)
                                 # dan kemudian gunakan watermark_style = 'B'
            # st.sidebar.info(f"Berjaya memuatkan fon Unicode: {UNICODE_FONT_FAMILY}") # Mesej debug
        except RuntimeError as e:
            st.warning(f"Gagal memuatkan fon Unicode '{FONT_REGULAR_PATH}': {e}. Menggunakan fon lalai '{DEFAULT_FALLBACK_FONT}'.")
            # current_font_family_for_content dan current_font_family_for_watermark kekal sebagai DEFAULT_FALLBACK_FONT
    else:
        st.warning(f"Fail fon Unicode '{FONT_REGULAR_PATH}' tidak ditemui. Menggunakan fon lalai '{DEFAULT_FALLBACK_FONT}'. Pastikan fail fon ada dalam direktori '{FONT_DIR}'.")
        # current_font_family_for_content dan current_font_family_for_watermark kekal sebagai DEFAULT_FALLBACK_FONT
    # --- PERUBAHAN TAMAT DI SINI ---

    if logo_path and os.path.exists(logo_path):
        try:
            img_width = 30 # mm
            page_width = pdf.w - 2 * pdf.l_margin
            x_logo = (page_width - img_width) / 2 + pdf.l_margin
            pdf.image(logo_path, x=x_logo, y=10, w=img_width)
            pdf.ln(25) # Ruang selepas logo
        except Exception as e:
            st.warning(f"Gagal menambah logo pada PDF: {e}. Pastikan fail imej sah dan format disokong oleh FPDF (PNG, JPG, GIF).")

    # Simpan kedudukan Y semasa sebelum tera air, jika ada logo
    y_before_watermark = pdf.get_y()

    if watermark_text:
        # --- PERUBAHAN: Gunakan fon yang telah ditentukan untuk tera air ---
        pdf.set_font(current_font_family_for_watermark, style=watermark_style, size=30)
        # --- TAMAT PERUBAHAN ---
        pdf.set_text_color(220, 220, 220) # Kelabu sangat cair
        text_w = pdf.get_string_width(watermark_text) # Perlu set fon dahulu sebelum get_string_width

        # Letakkan tera air di tengah halaman secara menegak dan mendatar
        page_center_x = pdf.w / 2
        page_center_y = pdf.h / 2
        pdf.set_xy(page_center_x - (text_w / 2), page_center_y - 5) # -5 untuk sedikit penyesuaian ketinggian teks

        pdf.cell(text_w, 10, watermark_text, 0, 0, 'C')
        pdf.set_text_color(0, 0, 0) # Reset warna teks

        # Reset kedudukan Y ke kedudukan sebelum tera air atau selepas logo
        # Ini penting supaya kandungan utama tidak bertindih dengan tera air atau logo
        pdf.set_xy(pdf.l_margin, y_before_watermark)
        if not (logo_path and os.path.exists(logo_path)): # Jika tiada logo, tambah sedikit ruang
            pdf.ln(5)


    # --- PERUBAHAN: Gunakan fon Unicode untuk kandungan utama ---
    pdf.set_font(current_font_family_for_content, size=12)
    # --- TAMAT PERUBAHAN ---

    for para_block in text_content.split("\n\n"):
        # FPDF multi_cell menangani pemecahan baris secara automatik jika teks terlalu panjang
        # Ia juga menghormati \n dalam string
        pdf.multi_cell(0, 10, para_block.strip())
        pdf.ln(5) # Ruang antara perenggan (blok)
    try:
        # Simpan ke fail tempatan. 'F' tidak diperlukan jika nama fail diberikan.
        pdf.output(filename)
        return True
    except Exception as e:
        st.error(f"Gagal menyimpan ke PDF: {e}")
        return False

def save_to_txt(text_content, filename='output.txt'):
    try:
        with open(filename, "w", encoding="utf-8") as f: f.write(text_content)
        return True
    except IOError as e: st.error(f"Gagal menyimpan ke Teks: {e}"); return False

def save_to_excel(chat_history, filename='chat_output.xlsx'):
    # Eksport semua mesej (user dan assistant) seperti dalam chatbot1
    data = [[msg["role"].capitalize(), msg["content"]] for msg in chat_history]
    df = pd.DataFrame(data, columns=["Role", "Message"])
    try: df.to_excel(filename, index=False, engine='openpyxl'); return True
    except Exception as e: st.error(f"Gagal menyimpan ke Excel: {e}"); return False

def save_to_pptx(chat_history, filename='chat_output.pptx', logo_path=LOGO_PATH):
    prs = Presentation()
    # Gunakan susun atur yang lebih fleksibel, contohnya 'Blank' atau 'Title and Content'
    # slide_layout = prs.slide_layouts[5] # Title Only
    slide_layout = prs.slide_layouts[6] # Blank, lebih fleksibel

    for msg in chat_history:
        slide = prs.slides.add_slide(slide_layout)

        if logo_path and os.path.exists(logo_path):
            try:
                # Logo di penjuru atas kiri
                pic = slide.shapes.add_picture(logo_path, PptxInches(0.2), PptxInches(0.2), height=PptxInches(0.75))
            except Exception as e:
                 st.warning(f"Gagal menambah logo pada PowerPoint: {e}. Pastikan fail imej sah.")

        # Kotak teks untuk peranan dan kandungan
        # Laraskan kedudukan dan saiz berdasarkan kehadiran logo
        left = PptxInches(0.5)
        top = PptxInches(1.0) if logo_path and os.path.exists(logo_path) else PptxInches(0.5)
        width = PptxInches(9.0)
        height = PptxInches(5.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        p_role = tf.add_paragraph()
        p_role.text = f"{msg['role'].capitalize()}:"
        p_role.font.bold = True
        p_role.font.size = PptxPt(18)
        p_role.font.name = 'Arial'

        p_content = tf.add_paragraph()
        p_content.text = msg['content']
        p_content.font.size = PptxPt(16)
        p_content.font.name = 'Arial'
        p_content.level = 1 # Inden sedikit untuk kandungan

    try: prs.save(filename); return True
    except IOError as e: st.error(f"Gagal menyimpan ke PowerPoint: {e}"); return False

# --- PENGURUSAN STATE STREAMLIT ---
def initialize_session_state(available_models_list):
    if "session_id" not in st.session_state:
        st.session_state.session_id = "new"
        st.session_state.chat_history = []
        st.session_state.current_filename_prefix = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    if "selected_ollama_model" not in st.session_state:
        if available_models_list and DEFAULT_OLLAMA_MODEL in available_models_list:
            st.session_state.selected_ollama_model = DEFAULT_OLLAMA_MODEL
        elif available_models_list:
            st.session_state.selected_ollama_model = available_models_list[0]
        else:
            st.session_state.selected_ollama_model = DEFAULT_OLLAMA_MODEL

    if "show_confirm_delete_all_button" not in st.session_state:
        st.session_state.show_confirm_delete_all_button = False
    
    if "chat_page_num" not in st.session_state:
        st.session_state.chat_page_num = 1

    # --- LOGIK BARU UNTUK KEY FILE UPLOADER ---
    if "uploader_key_counter" not in st.session_state:
        st.session_state.uploader_key_counter = 0
    # --- TAMAT LOGIK BARU ---


# --- KOMPONEN UI ---
def display_sidebar(available_models_list):
    st.sidebar.header("⚙️ Tetapan")
    if available_models_list:
        try:
            current_model_index = available_models_list.index(st.session_state.selected_ollama_model)
        except ValueError:
            current_model_index = 0
            if available_models_list: st.session_state.selected_ollama_model = available_models_list[0]
            else: st.session_state.selected_ollama_model = DEFAULT_OLLAMA_MODEL

        selected_model_ui = st.sidebar.selectbox(
            "Pilih Model AI:", options=available_models_list, index=current_model_index, key="model_selector_widget"
        )
        if selected_model_ui != st.session_state.selected_ollama_model:
            st.session_state.selected_ollama_model = selected_model_ui
            # st.rerun() # Tidak perlu rerun di sini, akan dikemas kini secara automatik
    else:
        st.sidebar.warning("Tiada model AI ditemui dari Ollama.")
        if "selected_ollama_model" not in st.session_state: # Pastikan ada nilai walaupun tiada model
             st.session_state.selected_ollama_model = DEFAULT_OLLAMA_MODEL

    st.sidebar.divider()
    st.sidebar.header("🕘 Sesi Perbualan")
    session_ids = load_all_session_ids()
    # Pastikan "new" (atau apa sahaja yang mewakili sesi baru) ada dalam session_ids jika itu ID semasa
    # Ini penting untuk selectbox mencari index yang betul
    current_session_for_select = st.session_state.session_id
    
    options = ["➕ Perbualan Baru"] + session_ids
    
    try:
        # Jika sesi semasa ialah "new", index adalah 0 ("➕ Perbualan Baru")
        # Jika tidak, cari index sesi semasa dalam senarai options
        current_session_index = options.index(current_session_for_select) if current_session_for_select != "new" else 0
    except ValueError: # Jika session_id semasa tidak ditemui (cth: fail dipadam secara manual)
        current_session_index = 0 # Default ke "➕ Perbualan Baru"
        st.session_state.session_id = "new"
        st.session_state.chat_history = []
        st.session_state.current_filename_prefix = datetime.now().strftime("%Y%m%d_%H%M%S")


    selected_session_id_ui = st.sidebar.selectbox(
        "Pilih atau mulakan sesi Perbualan:", options, index=current_session_index, key="session_selector_widget"
    )

    st.sidebar.divider()
    st.sidebar.subheader("🗑️ Urus Sesi")
    can_delete_current = st.session_state.session_id != "new" and st.session_state.session_id in session_ids
    if can_delete_current:
        if st.sidebar.button(f"Padam Sesi: {st.session_state.session_id}", key="delete_current_btn", type="secondary"):
            if delete_chat_session_file(st.session_state.session_id):
                st.session_state.session_id = "new"; st.session_state.chat_history = []
                st.session_state.current_filename_prefix = datetime.now().strftime("%Y%m%d_%H%M%S")
                st.session_state.show_confirm_delete_all_button = False
                st.session_state.chat_page_num = 1 # Reset paginasi
                st.rerun()

    if session_ids:
        if not st.session_state.show_confirm_delete_all_button:
            if st.sidebar.button("Padam Semua Sesi", key="ask_delete_all_btn"):
                st.session_state.show_confirm_delete_all_button = True
                st.rerun()
        
        if st.session_state.show_confirm_delete_all_button:
            st.sidebar.warning("ANDA PASTI MAHU MEMADAM SEMUA SESI? TINDAKAN INI TIDAK BOLEH DIBATALKAN.")
            col1, col2 = st.sidebar.columns(2)
            with col1:
                if st.button("YA, PADAM SEMUA", key="confirm_delete_all_btn", type="primary"):
                    if delete_all_chat_sessions():
                        st.session_state.session_id = "new"; st.session_state.chat_history = []
                        st.session_state.current_filename_prefix = datetime.now().strftime("%Y%m%d_%H%M%S")
                        st.session_state.show_confirm_delete_all_button = False
                        st.session_state.chat_page_num = 1 # Reset paginasi
                        st.rerun()
            with col2:
                if st.button("TIDAK, BATAL", key="cancel_delete_all_btn"):
                    st.session_state.show_confirm_delete_all_button = False
                    st.rerun()
    else:
        st.session_state.show_confirm_delete_all_button = False
    
    # st.sidebar.info(
    #     f"""
    #     **Nota:**
    #     Pastikan servis Ollama anda berjalan.
    #     Model yang tersedia akan disenaraikan di atas.
    #     URL Ollama: `{OLLAMA_BASE_URL}`
    #     Logo: `{LOGO_PATH if os.path.exists(LOGO_PATH) else "Tidak ditemui"}`
    #     Tera Air: `{WATERMARK_TEXT}`
    #     """
    # )
    return selected_session_id_ui

def handle_session_logic(selected_session_id_from_ui):
    # Logik ini dijalankan apabila pilihan sesi di sidebar berubah
    if selected_session_id_from_ui == "➕ Perbualan Baru":
        if st.session_state.session_id != "new": # Jika bertukar DARI sesi sedia ada KE baru
            st.session_state.session_id = "new"
            st.session_state.chat_history = []
            st.session_state.current_filename_prefix = datetime.now().strftime("%Y%m%d_%H%M%S")
            st.session_state.chat_page_num = 1
            # st.rerun() # Rerun akan berlaku secara semula jadi jika widget berubah
    elif st.session_state.session_id != selected_session_id_from_ui: # Jika bertukar KE sesi sedia ada YANG LAIN
        st.session_state.chat_history = load_chat_session(selected_session_id_from_ui)
        st.session_state.session_id = selected_session_id_from_ui
        st.session_state.current_filename_prefix = selected_session_id_from_ui # Gunakan ID sesi sebagai prefix
        st.session_state.chat_page_num = 1
        # st.rerun()

def display_chat_messages_paginated():
    st.subheader("📜 Perbualan")
    if not st.session_state.chat_history:
        st.info("Mulakan perbualan dengan menaip di bawah atau muat naik fail.")
        return

    page_size = 10 # Bilangan mesej setiap halaman
    total_messages = len(st.session_state.chat_history)
    
    # Kira max_page dengan betul, pastikan sekurang-kurangnya 1 halaman
    max_page = (total_messages + page_size - 1) // page_size if total_messages > 0 else 1
    
    # Pastikan chat_page_num berada dalam lingkungan yang sah
    if st.session_state.chat_page_num > max_page:
        st.session_state.chat_page_num = max_page
    if st.session_state.chat_page_num < 1:
        st.session_state.chat_page_num = 1

    # Hanya tunjukkan slider jika lebih dari satu halaman
    if max_page > 1:
        page_num_ui = st.slider(
            "Halaman Perbualan:", 
            min_value=1, 
            max_value=max_page, 
            value=st.session_state.chat_page_num, 
            key="chat_page_slider"
        )
        if page_num_ui != st.session_state.chat_page_num:
            st.session_state.chat_page_num = page_num_ui
            # st.rerun() # Tidak perlu rerun, widget akan trigger
    else:
        st.session_state.chat_page_num = 1 # Jika hanya satu halaman, pastikan ia adalah halaman 1

    # Kira mesej untuk dipaparkan berdasarkan halaman semasa
    # Papar mesej terbaru dahulu (songsangkan senarai untuk paparan, tetapi simpan dalam susunan asal)
    reversed_history = st.session_state.chat_history[::-1]
    start_index = (st.session_state.chat_page_num - 1) * page_size
    end_index = start_index + page_size
    
    # Papar mesej dalam susunan kronologi untuk halaman semasa
    messages_to_display = reversed_history[start_index:end_index][::-1] 

    for msg in messages_to_display:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])
            if msg["role"] == "assistant" and "time_taken" in msg and msg["time_taken"] is not None:
                st.caption(f"Dijana dalam {msg['time_taken']:.2f} saat")

def display_export_options():
    st.divider()
    st.subheader("📤 Eksport Perbualan")
    if not st.session_state.chat_history:
        st.info("Tiada perbualan untuk dieksport.")
        return

    col_export1, col_export2 = st.columns(2)
    with col_export1:
        export_content_choice = st.radio(
            "Kandungan untuk dieksport (Teks, Word, PDF, Excel, PPTX):", # Kemas kini label
            ["Pembantu Sahaja", "Pengguna Sahaja", "Keseluruhan Perbualan"],
            index=2, key="export_content_radio"
        )
    with col_export2:
        export_format_choice = st.selectbox("Format eksport:", [
            "Pilih format", "Word (.docx)", "Teks (.txt)", "PDF (.pdf)",
            "Excel (.xlsx)", "PowerPoint (.pptx)"
        ], key="export_format_select")

    custom_filename_prefix_ui = st.text_input(
        "Nama fail awalan (tanpa sambungan):",
        st.session_state.current_filename_prefix,
        key="filename_prefix_input"
    )
    if custom_filename_prefix_ui != st.session_state.current_filename_prefix:
        st.session_state.current_filename_prefix = custom_filename_prefix_ui

    if st.button("Eksport", key="export_main_button"):
        if export_format_choice == "Pilih format":
            st.warning("Sila pilih format eksport yang sah."); return
        
        filename_base = st.session_state.current_filename_prefix
        include_user = "Pengguna" in export_content_choice or "Keseluruhan" in export_content_choice
        include_assistant = "Pembantu" in export_content_choice or "Keseluruhan" in export_content_choice
        
        text_for_common_formats = format_conversation_text(st.session_state.chat_history, include_user, include_assistant)
        
        # Tapis sejarah untuk Excel dan PowerPoint berdasarkan pilihan radio
        history_for_excel_pptx = [
            msg for msg in st.session_state.chat_history 
            if (include_user and msg["role"] == "user") or \
               (include_assistant and msg["role"] == "assistant")
        ]
        # Pastikan history_for_excel_pptx tidak kosong jika pengguna memilih untuk eksport hanya satu peranan
        # dan peranan itu tiada dalam sejarah. Fungsi eksport mungkin gagal.
        if not history_for_excel_pptx and (export_format_choice == "Excel (.xlsx)" or export_format_choice == "PowerPoint (.pptx)"):
            st.warning(f"Tiada mesej '{export_content_choice.lower().replace(' keseluruhan perbualan', '')}' ditemui untuk dieksport ke {export_format_choice}.")
            return

        success, exported_filename = False, ""
        
        # Fungsi save_to_word dan save_to_pdf menggunakan LOGO_PATH dan WATERMARK_TEXT secara lalai dari pembolehubah global.
        # Fungsi save_to_pptx menggunakan LOGO_PATH secara lalai.
        actions = {
            "Word (.docx)": (save_to_word, text_for_common_formats, f"{filename_base}.docx"),
            "Teks (.txt)": (save_to_txt, text_for_common_formats, f"{filename_base}.txt"),
            "PDF (.pdf)": (save_to_pdf, text_for_common_formats, f"{filename_base}.pdf"),
            "Excel (.xlsx)": (save_to_excel, history_for_excel_pptx, f"{filename_base}.xlsx"),
            "PowerPoint (.pptx)": (save_to_pptx, history_for_excel_pptx, f"{filename_base}.pptx")
        }

        if export_format_choice in actions:
            func, data_to_export, fname = actions[export_format_choice]
            # Semak jika data untuk dieksport kosong (terutamanya untuk Excel/PPTX selepas penapisan)
            if not data_to_export: # Jika data_to_export adalah senarai kosong atau string kosong
                 st.warning(f"Tiada kandungan untuk dieksport ke {export_format_choice} berdasarkan pilihan anda.")
                 return

            success = func(data_to_export, fname)
            exported_filename = fname
        
        if success and exported_filename:
            st.success(f"Fail disimpan: {exported_filename}")
            try:
                with open(exported_filename, "rb") as f_download:
                    st.download_button(
                        "📥 Muat Turun Fail", data=f_download, file_name=exported_filename, 
                        key=f"download_btn_{exported_filename.replace('.', '_')}_{time.time()}"
                    )
            except FileNotFoundError: st.error(f"Gagal mencari fail {exported_filename} untuk dimuat turun.")
            except Exception as e: st.error(f"Ralat semasa menyediakan muat turun: {e}")

# --- FUNGSI UTAMA APLIKASI ---
def main():
    st.set_page_config(page_title="DFK Stembot", layout="wide", initial_sidebar_state="expanded", page_icon="🤖")
    st.title("🤖 DFK Stembot")

    available_ollama_models = get_ollama_models_cached()
    if not available_ollama_models:
        st.error("Tidak dapat memuatkan senarai model dari Ollama. Pastikan Ollama berjalan dan mempunyai model. Aplikasi mungkin tidak berfungsi dengan betul.")
    
    initialize_session_state(available_ollama_models) # current_filename_prefix diinisialisasi di sini untuk sesi "new"

    st.caption(f"Model semasa: **{st.session_state.selected_ollama_model}**")
    
    selected_session_id_from_ui = display_sidebar(available_ollama_models)
    handle_session_logic(selected_session_id_from_ui) # Mengendalikan pemuatan sesi sedia ada atau reset ke "new"
    
    # --- Bahagian Muat Naik Fail ---
    st.sidebar.divider()
    st.sidebar.header("📎 Muat Naik Fail")
    
    uploader_key = f"file_uploader_{st.session_state.uploader_key_counter}"
    uploaded_file = st.sidebar.file_uploader(
        "Muat naik imej, PDF, DOCX, atau TXT untuk diproses:", 
        type=['png', 'jpg', 'jpeg', 'gif', 'pdf', 'txt', 'docx'],
        key=uploader_key
    )

    if uploaded_file is not None:
        with st.spinner(f"Memproses fail '{uploaded_file.name}'..."):
            extracted_text = extract_text_from_file(uploaded_file)
        
        if extracted_text:
            st.info(f"Teks diekstrak dari '{uploaded_file.name}'. Anda boleh bertanya mengenainya atau ia akan disertakan dalam konteks seterusnya.")
            file_content_message = f"Kandungan dari fail '{uploaded_file.name}':\n\n{extracted_text}"
            
            st.session_state.chat_history.append({"role": "user", "content": file_content_message})
            
            with st.spinner(f"{st.session_state.selected_ollama_model.split(':')[0].capitalize()} sedang memproses kandungan fail..."):
                assistant_response, gen_time = query_ollama_non_stream(
                    file_content_message,
                    st.session_state.chat_history, 
                    st.session_state.selected_ollama_model
                )

            st.session_state.chat_history.append({
                "role": "assistant", 
                "content": assistant_response,
                "time_taken": gen_time
            })
            
            # --- LOGIK PENYIMPANAN DIPERBAIKI ---
            if st.session_state.session_id == "new":
                # Ini adalah mesej pertama dalam sesi baru.
                # Gunakan current_filename_prefix (yang sepatutnya cap masa) sebagai ID sesi baru.
                st.session_state.session_id = st.session_state.current_filename_prefix
                # Selepas ini, session_id tidak lagi "new" untuk interaksi seterusnya dalam sesi ini.
            
            # Simpan sesi (sama ada sesi baru yang IDnya baru ditetapkan, atau sesi sedia ada yang dikemas kini)
            save_chat_session(st.session_state.session_id, st.session_state.chat_history)
            # --- TAMAT LOGIK PENYIMPANAN DIPERBAIKI ---
        
        elif extracted_text is None: 
            # Mesej ralat/amaran sudah dipaparkan oleh extract_text_from_file
            pass
        else: # extracted_text adalah string kosong
            st.warning(f"Tiada teks dapat diekstrak dari fail '{uploaded_file.name}'.")

        st.session_state.uploader_key_counter += 1
        st.rerun() # Rerun diperlukan untuk memaparkan mesej baru dan mengosongkan pemuat naik fail

    display_chat_messages_paginated()

    friendly_model_name = st.session_state.selected_ollama_model.split(':')[0].capitalize()
    user_input = st.chat_input(f"Taip mesej anda kepada {friendly_model_name}...")

    if user_input:
        st.session_state.chat_history.append({"role": "user", "content": user_input})
        
        with st.spinner(f"{friendly_model_name} sedang menaip..."): 
            assistant_response_text, generation_time = query_ollama_non_stream(
                user_input, 
                st.session_state.chat_history, 
                st.session_state.selected_ollama_model
            )
        
        st.session_state.chat_history.append({
            "role": "assistant", 
            "content": assistant_response_text,
            "time_taken": generation_time
        })

        # --- LOGIK PENYIMPANAN DIPERBAIKI ---
        if st.session_state.session_id == "new":
            # Ini adalah mesej pertama dalam sesi baru.
            # Gunakan current_filename_prefix (yang sepatutnya cap masa) sebagai ID sesi baru.
            st.session_state.session_id = st.session_state.current_filename_prefix
            # Selepas ini, session_id tidak lagi "new" untuk interaksi seterusnya dalam sesi ini.
        
        # Simpan sesi (sama ada sesi baru yang IDnya baru ditetapkan, atau sesi sedia ada yang dikemas kini)
        save_chat_session(st.session_state.session_id, st.session_state.chat_history)
        # --- TAMAT LOGIK PENYIMPANAN DIPERBAIKI ---
        
        total_messages = len(st.session_state.chat_history)
        page_size = 10 
        st.session_state.chat_page_num = (total_messages + page_size - 1) // page_size if total_messages > 0 else 1
        
        st.rerun() # Rerun diperlukan untuk memaparkan mesej baru

    display_export_options()


if __name__ == "__main__":
    main()
        # Cadangan untuk fail requirements.txt:
    #
    # streamlit>=1.17.0
    # requests
    # python-docx
    # fpdf2
    # pandas
    # openpyxl
    # python-pptx
    # Pillow
    # pytesseract
    # PyMuPDF
    #
    # Pasang menggunakan: pip install -r requirements.txt
    # Juga, pastikan Tesseract OCR dipasang pada sistem anda.
