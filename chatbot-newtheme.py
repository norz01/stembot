import streamlit as st
import requests
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from fpdf import FPDF
import pandas as pd
from datetime import datetime
import os
import json
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
import time
from PIL import Image
import pytesseract
import fitz
import bcrypt

# --- KONFIGURASI ---
HISTORY_DIR = "chat_sessions"  # Direktori utama untuk semua sesi pengguna
UPLOAD_DIR = "uploaded_files"
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
DEFAULT_OLLAMA_MODEL = os.getenv("DEFAULT_OLLAMA_MODEL", "STEMBot-4B") # Pastikan model ini wujud
LOGO_PATH = os.getenv("ikm_logo", "ikm_logo.jpg") # Guna pembolehubah ini
WATERMARK_TEXT = os.getenv("CHATBOT_WATERMARK_TEXT", "IKM Besut")
USERS_DIR = "user_data"
USERS_FILE = os.path.join(USERS_DIR, "users.json")

# Pastikan direktori wujud
os.makedirs(HISTORY_DIR, exist_ok=True) # Hanya direktori utama
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(USERS_DIR, exist_ok=True)

# --- FUNGSI PENGURUSAN AKAUN (Tidak Berubah) ---
def load_users():
    if not os.path.exists(USERS_FILE):
        with open(USERS_FILE, "w") as f:
            json.dump({}, f)
    with open(USERS_FILE, "r") as f:
        return json.load(f)

def save_users(users):
    with open(USERS_FILE, "w") as f:
        json.dump(users, f, indent=2)

def hash_password(password):
    return bcrypt.hashpw(password.encode(), bcrypt.gensalt()).decode()

def verify_password(password, hashed):
    return bcrypt.checkpw(password.encode(), hashed.encode())

# --- HALAMAN LOGIN & PENDAFTARAN (Tidak Berubah) ---
def login_page():
    st.title("🔐 Log Masuk")
    username = st.text_input("Nama Pengguna", key="login_username")
    password = st.text_input("Kata Laluan", type="password", key="login_password")
    
    if st.button("Log Masuk", type="primary", use_container_width=True):
        users = load_users()
        if username in users and verify_password(password, users[username]["password"]):
            st.session_state.authenticated = True
            st.session_state.username = username
            st.success("Berjaya log masuk!")
            st.rerun()
        else:
            st.error("Nama pengguna atau kata laluan salah.")

def register_page():
    st.title("📝 Daftar Akaun")
    username = st.text_input("Nama Pengguna Baru", key="register_username")
    password = st.text_input("Kata Laluan", type="password", key="register_password")
    confirm_password = st.text_input("Sahkan Kata Laluan", type="password", key="confirm_password")
    
    if st.button("Daftar", type="primary", use_container_width=True):
        if not username or not password:
            st.warning("Sila lengkapkan semua medan.")
            return
        if password != confirm_password:
            st.error("Kata laluan tidak sepadan.")
            return
        users = load_users()
        if username in users:
            st.error("Nama pengguna telah wujud.")
            return
        users[username] = {
            "password": hash_password(password),
            "created_at": datetime.now().isoformat()
        }
        save_users(users)
        st.success("Akaun berjaya didaftarkan! Sila log masuk.")
        # st.session_state.show_register = False # Tidak perlu jika guna tabs
        st.rerun()

def authentication_ui():
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
    
    if not st.session_state.authenticated:
        st.markdown("## Selamat Datang ke DFK Stembot")
        tab1, tab2 = st.tabs(["Log Masuk", "Daftar Akaun"])
        with tab1:
            login_page()
        with tab2:
            register_page()
        return False
    return True

# --- FUNGSI HELPER ---
@st.cache_data(ttl=300)
def get_ollama_models_cached():
    try:
        response = requests.get(f'{OLLAMA_BASE_URL}/api/tags', timeout=10)
        response.raise_for_status()
        models_data = response.json().get('models', [])
        if not models_data:
            return []
        return sorted([model['name'] for model in models_data])
    except requests.exceptions.RequestException:
        # st.error("Gagal menyambung ke Ollama untuk mendapatkan senarai model.")
        return []
    except Exception: # Tangkap semua yang lain
        # st.error("Ralat tidak dijangka semasa mendapatkan senarai model.")
        return []

def query_ollama_non_stream(prompt, chat_history, selected_model):
    messages_for_api = [{"role": msg["role"], "content": msg["content"]} for msg in chat_history]
    is_prompt_already_last_user_message = False
    if messages_for_api and messages_for_api[-1]["role"] == "user" and messages_for_api[-1]["content"] == prompt:
        is_prompt_already_last_user_message = True
    if not is_prompt_already_last_user_message:
         messages_for_api.append({"role": "user", "content": prompt})

    start_time = time.time()
    thinking_process = "" 
    assistant_reply = "Maaf, saya tidak dapat respons yang betul." 
    try:
        payload = {'model': selected_model, 'messages': messages_for_api, 'stream': False}
        response = requests.post(f'{OLLAMA_BASE_URL}/api/chat', json=payload, timeout=600) 
        response.raise_for_status()
        full_response_data = response.json()
        raw_assistant_reply = full_response_data.get('message', {}).get('content')
        if raw_assistant_reply is None:
            raw_assistant_reply = "Maaf, respons dari model tidak mengandungi kandungan."
            assistant_reply = raw_assistant_reply
        else:
            assistant_reply = raw_assistant_reply
        thinking_start_tag = "<think>" # Tag pemikiran anda
        thinking_end_tag = "</think>"
        if thinking_start_tag in raw_assistant_reply and thinking_end_tag in raw_assistant_reply:
            try:
                start_index = raw_assistant_reply.find(thinking_start_tag)
                end_index = raw_assistant_reply.find(thinking_end_tag)
                if 0 <= start_index < end_index:
                    thinking_process = raw_assistant_reply[start_index + len(thinking_start_tag):end_index].strip()
                    text_after_thinking = raw_assistant_reply[end_index + len(thinking_end_tag):].strip()
                    text_before_thinking = raw_assistant_reply[:start_index].strip()
                    if text_after_thinking:
                        assistant_reply = text_after_thinking
                    elif text_before_thinking:
                        assistant_reply = text_before_thinking
                    elif thinking_process: 
                        assistant_reply = "" 
            except Exception as e_parse:
                st.warning(f"Ralat kecil semasa memproses tag pemikiran: {e_parse}. Menggunakan respons penuh.")
                assistant_reply = raw_assistant_reply
                thinking_process = ""
        end_time = time.time()
        processing_time = end_time - start_time
        return assistant_reply, thinking_process, processing_time
    except requests.exceptions.HTTPError as http_err:
        end_time = time.time(); processing_time = end_time - start_time
        st.error(f"Ralat HTTP dari Ollama: {http_err} (selepas {processing_time:.2f}s)")
        try:
            error_details = response.json().get("error", "Tiada butiran ralat tambahan.")
            st.error(f"Butiran dari Ollama: {error_details}")
            return f"Maaf, berlaku ralat HTTP semasa menghubungi Ollama: {error_details}", "", processing_time
        except:
             return f"Maaf, berlaku ralat HTTP semasa menghubungi Ollama ({http_err.response.status_code}).", "", processing_time
    except requests.exceptions.Timeout:
        end_time = time.time(); processing_time = end_time - start_time
        st.error(f"Gagal mendapatkan respons: Permintaan ke Ollama tamat masa selepas {processing_time:.2f}s.")
        return "Maaf, permintaan tamat masa.", "", processing_time
    except requests.exceptions.RequestException as e:
        end_time = time.time(); processing_time = end_time - start_time
        st.error(f"Masalah menyambung ke Ollama: {e} (selepas {processing_time:.2f}s)")
        return "Maaf, berlaku masalah semasa menghubungi Ollama.", "", processing_time
    except json.JSONDecodeError:
        end_time = time.time(); processing_time = end_time - start_time
        st.error(f"Format respons tidak dijangka (bukan JSON) dari Ollama (selepas {processing_time:.2f}s).")
        return "Maaf, format respons dari Ollama tidak seperti yang dijangkakan.", "", processing_time
    except KeyError:
        end_time = time.time(); processing_time = end_time - start_time
        st.error(f"Struktur data respons tidak dijangka dari Ollama (selepas {processing_time:.2f}s).")
        return "Maaf, struktur data respons dari Ollama tidak seperti yang dijangkakan.", "", processing_time
    except Exception as e:
        end_time = time.time(); processing_time = end_time - start_time
        st.error(f"Ralat tidak dijangka dalam query_ollama_non_stream: {e} (selepas {processing_time:.2f}s)")
        return "Maaf, ralat tidak dijangka berlaku semasa memproses permintaan.", "", processing_time

# --- PERUBAHAN PADA FUNGSI PENGURUSAN SESI ---
def get_user_history_dir(username):
    """Mendapatkan direktori sejarah untuk pengguna tertentu."""
    user_dir = os.path.join(HISTORY_DIR, username)
    os.makedirs(user_dir, exist_ok=True) # Pastikan direktori pengguna wujud
    return user_dir

def save_chat_session(username, session_id, history):
    """Menyimpan sesi perbualan untuk pengguna tertentu."""
    user_history_dir = get_user_history_dir(username)
    filepath = os.path.join(user_history_dir, f"{session_id}.json")
    try:
        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(history, f, indent=2)
    except IOError as e:
        st.error(f"Gagal menyimpan sesi Perbualan '{session_id}' untuk pengguna '{username}': {e}")

def load_chat_session(username, session_id):
    """Memuatkan sesi perbualan untuk pengguna tertentu."""
    user_history_dir = get_user_history_dir(username)
    filepath = os.path.join(user_history_dir, f"{session_id}.json")
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return json.load(f)
    except FileNotFoundError:
        return []
    except (json.JSONDecodeError, IOError) as e:
        st.error(f"Gagal memuatkan sesi Perbualan '{session_id}' untuk pengguna '{username}': {e}")
        return []

def load_all_session_ids(username):
    """Memuatkan semua ID sesi untuk pengguna tertentu."""
    user_history_dir = get_user_history_dir(username)
    try:
        if not os.path.exists(user_history_dir): # Jika pengguna belum ada sesi
            return []
        files = [f.replace(".json", "") for f in os.listdir(user_history_dir) if f.endswith(".json")]
        def sort_key(filename):
            try:
                parts = filename.split('_')
                if len(parts) >= 2:
                    return datetime.strptime(f"{parts[0]}_{parts[1]}", "%Y%m%d_%H%M%S")
            except (ValueError, IndexError): pass
            return datetime.min 
        return sorted(files, key=sort_key, reverse=True)
    except OSError as e:
        st.error(f"Gagal membaca direktori sesi untuk pengguna '{username}': {e}")
        return []

def delete_chat_session_file(username, session_id):
    """Memadam fail sesi perbualan untuk pengguna tertentu."""
    user_history_dir = get_user_history_dir(username)
    filepath = os.path.join(user_history_dir, f"{session_id}.json")
    try:
        if os.path.exists(filepath):
            os.remove(filepath)
            st.success(f"Sesi Perbualan '{session_id}' berjaya dipadam.")
            return True
        else:
            st.warning(f"Fail sesi Perbualan '{session_id}' tidak ditemui untuk dipadam.")
            return False
    except OSError as e:
        st.error(f"Gagal memadam sesi Perbualan '{session_id}': {e}")
        return False

def delete_all_chat_sessions(username):
    """Memadam semua sesi perbualan untuk pengguna tertentu."""
    user_history_dir = get_user_history_dir(username)
    deleted_count = 0
    errors = []
    try:
        if not os.path.exists(user_history_dir):
            st.info("Tiada sesi Perbualan ditemui untuk pengguna ini.")
            return True # Dianggap berjaya kerana tiada apa untuk dipadam

        for filename in os.listdir(user_history_dir):
            if filename.endswith(".json"):
                filepath = os.path.join(user_history_dir, filename)
                try: 
                    os.remove(filepath)
                    deleted_count += 1
                except OSError as e: 
                    errors.append(f"Gagal memadam {filename}: {e}")
        if errors:
            for error in errors: st.error(error)
        if deleted_count > 0: 
            st.success(f"{deleted_count} sesi Perbualan untuk pengguna '{username}' berjaya dipadam.")
        else: 
            st.info(f"Tiada sesi Perbualan ditemui untuk pengguna '{username}' untuk dipadam.")
        return True
    except OSError as e:
        st.error(f"Gagal mengakses direktori sesi untuk pengguna '{username}': {e}")
        return False
# --- TAMAT PERUBAHAN FUNGSI PENGURUSAN SESI ---

# ... (Fungsi extract_text_from_file, format_conversation_text, save_to_word, dll. tidak berubah) ...
# Pastikan semua fungsi ini masih ada. Saya akan sertakan format_conversation_text dan fungsi eksport yang diubah suai sedikit.

def extract_text_from_file(uploaded_file_obj):
    # ... (Kod sedia ada anda, tidak berubah) ...
    extracted_text = ""
    filename = uploaded_file_obj.name
    file_bytes = uploaded_file_obj.getvalue() 
    temp_docx_path = None 
    try:
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
            image = Image.open(uploaded_file_obj) 
            extracted_text = pytesseract.image_to_string(image)
            if not extracted_text.strip():
                 st.info(f"Tiada teks dapat diekstrak dari imej '{filename}' menggunakan OCR.")
        elif filename.lower().endswith(".txt"):
            extracted_text = file_bytes.decode('utf-8', errors='ignore')
        elif filename.lower().endswith(".docx"):
            temp_docx_path = os.path.join(UPLOAD_DIR, f"temp_{filename}") 
            try:
                with open(temp_docx_path, "wb") as f:
                    f.write(file_bytes)
                doc = Document(temp_docx_path)
                extracted_text = "\n".join([para.text for para in doc.paragraphs])
            except Exception as e_docx:
                st.error(f"Ralat semasa memproses fail DOCX '{filename}': {e_docx}")
            finally:
                if temp_docx_path and os.path.exists(temp_docx_path): 
                    try:
                        os.remove(temp_docx_path)
                    except OSError as e_remove:
                        st.warning(f"Gagal memadam fail sementara DOCX '{temp_docx_path}': {e_remove}")
        elif filename.lower().endswith(".pdf"):
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            for page in doc:
                extracted_text += page.get_text()
            doc.close()
        else:
            st.warning(f"Jenis fail '{filename}' tidak disokong untuk ekstraksi teks.")
            return None 
        return extracted_text.strip() if isinstance(extracted_text, str) else None
    except Exception as e:
        st.error(f"Ralat umum semasa memproses fail '{filename}': {e}")
        if temp_docx_path and os.path.exists(temp_docx_path): 
            try:
                os.remove(temp_docx_path)
            except OSError as e_remove_outer:
                st.warning(f"Gagal memadam fail sementara DOCX (luar) '{temp_docx_path}': {e_remove_outer}")
        return None

def format_conversation_text(chat_history, include_user=True, include_assistant=True):
    # ... (Kod sedia ada anda, tidak berubah) ...
    lines = []
    for msg in chat_history:
        role_display = msg["role"].capitalize()
        content_display = msg.get("content", "").strip()
        thinking_display = msg.get("thinking_process", "").strip()

        if (msg["role"] == "user" and include_user):
            lines.append(f"{role_display}: {content_display}")
        elif (msg["role"] == "assistant" and include_assistant):
            main_line = f"{role_display}: {content_display if content_display else '(Tiada jawapan utama)'}"
            lines.append(main_line)
            if thinking_display:
                lines.append(f"  Proses Pemikiran AI:\n  ---------------------\n{thinking_display}\n  ---------------------")
    return "\n\n".join(lines)

def save_to_word(text_content, filename='output.docx', logo_path=LOGO_PATH, watermark_text=WATERMARK_TEXT):
    # ... (Kod sedia ada anda, tidak berubah) ...
    doc = Document()
    if logo_path and os.path.exists(logo_path):
        try:
            paragraph = doc.add_paragraph()
            run = paragraph.add_run()
            run.add_picture(logo_path, width=DocxInches(2.0)) 
            paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            doc.add_paragraph() 
        except Exception as e:
            st.warning(f"Gagal menambah logo pada Word: {e}. Pastikan fail imej sah.")
    if watermark_text:
        watermark_para = doc.add_paragraph()
        run = watermark_para.add_run(watermark_text)
        font = run.font
        font.size = DocxPt(36) 
        font.color.rgb = DocxRGBColor(192, 192, 192)  
        font.bold = True
        watermark_para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        doc.add_paragraph()
    for para_block in text_content.split("\n\n"):
        doc.add_paragraph(para_block.strip())
    try: doc.save(filename); return True
    except IOError as e: st.error(f"Gagal menyimpan ke Word: {e}"); return False

def save_to_pdf(text_content, filename='output.pdf', logo_path=LOGO_PATH, watermark_text=WATERMARK_TEXT):
    # ... (Kod sedia ada anda, tidak berubah) ...
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    FONT_DIR = "fonts"
    FONT_REGULAR_FILENAME = "DejaVuSans.ttf"
    FONT_REGULAR_PATH = os.path.join(FONT_DIR, FONT_REGULAR_FILENAME)
    UNICODE_FONT_FAMILY = "DejaVuSans"
    DEFAULT_FALLBACK_FONT = "Arial"
    current_font_family_for_content = DEFAULT_FALLBACK_FONT
    current_font_family_for_watermark = DEFAULT_FALLBACK_FONT
    watermark_style = 'B'
    if os.path.exists(FONT_REGULAR_PATH):
        try:
            pdf.add_font(UNICODE_FONT_FAMILY, '', FONT_REGULAR_PATH, uni=True)
            current_font_family_for_content = UNICODE_FONT_FAMILY
            current_font_family_for_watermark = UNICODE_FONT_FAMILY
            watermark_style = '' 
        except RuntimeError as e:
            st.warning(f"Gagal memuatkan fon Unicode '{FONT_REGULAR_PATH}': {e}. Menggunakan fon lalai '{DEFAULT_FALLBACK_FONT}'.")
    else:
        st.warning(f"Fail fon Unicode '{FONT_REGULAR_PATH}' tidak ditemui. Menggunakan fon lalai '{DEFAULT_FALLBACK_FONT}'. Pastikan fail fon ada dalam direktori '{FONT_DIR}'.")
    if logo_path and os.path.exists(logo_path):
        try:
            img_width = 30 
            page_width = pdf.w - 2 * pdf.l_margin
            x_logo = (page_width - img_width) / 2 + pdf.l_margin
            pdf.image(logo_path, x=x_logo, y=10, w=img_width)
            pdf.ln(25) 
        except Exception as e:
            st.warning(f"Gagal menambah logo pada PDF: {e}. Pastikan fail imej sah dan format disokong oleh FPDF (PNG, JPG, GIF).")
    y_before_watermark = pdf.get_y()
    if watermark_text:
        pdf.set_font(current_font_family_for_watermark, style=watermark_style, size=30)
        pdf.set_text_color(220, 220, 220) 
        text_w = pdf.get_string_width(watermark_text) 
        page_center_x = pdf.w / 2
        page_center_y = pdf.h / 2
        pdf.set_xy(page_center_x - (text_w / 2), page_center_y - 5) 
        pdf.cell(text_w, 10, watermark_text, 0, 0, 'C')
        pdf.set_text_color(0, 0, 0) 
        pdf.set_xy(pdf.l_margin, y_before_watermark)
        if not (logo_path and os.path.exists(logo_path)): 
            pdf.ln(5)
    pdf.set_font(current_font_family_for_content, size=12)
    for para_block in text_content.split("\n\n"):
        pdf.multi_cell(0, 10, para_block.strip())
        pdf.ln(5) 
    try:
        pdf.output(filename)
        return True
    except Exception as e:
        st.error(f"Gagal menyimpan ke PDF: {e}")
        return False

def save_to_txt(text_content, filename='output.txt'):
    # ... (Kod sedia ada anda, tidak berubah) ...
    try:
        with open(filename, "w", encoding="utf-8") as f: f.write(text_content)
        return True
    except IOError as e: st.error(f"Gagal menyimpan ke Teks: {e}"); return False

def save_to_excel(chat_history, filename='chat_output.xlsx'):
    # ... (Kod sedia ada anda, tidak berubah) ...
    data = []
    for msg in chat_history:
        role = msg["role"].capitalize()
        content = msg.get("content", "")
        thinking = msg.get("thinking_process", "")
        if thinking:
            data.append([role, content, thinking])
        else:
            data.append([role, content, ""]) 
            
    df = pd.DataFrame(data, columns=["Role", "Message", "Thinking Process"])
    try: 
        df.to_excel(filename, index=False, engine='openpyxl')
        return True
    except Exception as e: 
        st.error(f"Gagal menyimpan ke Excel: {e}")
        return False

def save_to_pptx(chat_history, filename='chat_output.pptx', logo_path=LOGO_PATH):
    # ... (Kod sedia ada anda, tidak berubah) ...
    prs = Presentation()
    slide_layout = prs.slide_layouts[6] 
    for msg in chat_history:
        slide = prs.slides.add_slide(slide_layout)
        if logo_path and os.path.exists(logo_path):
            try:
                pic = slide.shapes.add_picture(logo_path, PptxInches(0.2), PptxInches(0.2), height=PptxInches(0.75))
            except Exception as e:
                 st.warning(f"Gagal menambah logo pada PowerPoint: {e}. Pastikan fail imej sah.")
        left = PptxInches(0.5)
        top = PptxInches(1.0) if logo_path and os.path.exists(logo_path) else PptxInches(0.5)
        width = PptxInches(9.0)
        height = PptxInches(5.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        p_role = tf.add_paragraph()
        p_role.text = f"{msg['role'].capitalize()}:"
        p_role.font.bold = True
        p_role.font.size = PptxPt(18)
        p_role.font.name = 'Arial'
        
        content_text = msg.get("content", "")
        thinking_text = msg.get("thinking_process", "")

        p_content = tf.add_paragraph()
        p_content.text = content_text
        p_content.font.size = PptxPt(16)
        p_content.font.name = 'Arial'
        p_content.level = 1 

        if thinking_text:
            p_thinking_header = tf.add_paragraph()
            run_thinking_header = p_thinking_header.add_run()
            run_thinking_header.text = "Proses Pemikiran AI:"
            run_thinking_header.font.italic = True
            run_thinking_header.font.size = PptxPt(14)
            
            p_thinking_content = tf.add_paragraph()
            p_thinking_content.text = thinking_text
            p_thinking_content.font.size = PptxPt(12)
            p_thinking_content.level = 2 
            
    try: prs.save(filename); return True
    except IOError as e: st.error(f"Gagal menyimpan ke PowerPoint: {e}"); return False

def initialize_session_state(available_models_list):
    # ... (Kod sedia ada anda, tidak berubah) ...
    if "session_id" not in st.session_state:
        st.session_state.session_id = "new"
        st.session_state.chat_history = []
        st.session_state.current_filename_prefix = datetime.now().strftime("%Y%m%d_%H%M%S")
    if "selected_ollama_model" not in st.session_state:
        if available_models_list and DEFAULT_OLLAMA_MODEL in available_models_list:
            st.session_state.selected_ollama_model = DEFAULT_OLLAMA_MODEL
        elif available_models_list:
            st.session_state.selected_ollama_model = available_models_list[0]
        else:
            st.session_state.selected_ollama_model = DEFAULT_OLLAMA_MODEL 
    if "show_confirm_delete_all_button" not in st.session_state:
        st.session_state.show_confirm_delete_all_button = False
    if "chat_page_num" not in st.session_state:
        st.session_state.chat_page_num = 1
    if "uploader_key_counter" not in st.session_state:
        st.session_state.uploader_key_counter = 0

# --- PERUBAHAN PADA FUNGSI UI ---
def display_sidebar(available_models_list, username): # Tambah username
    with st.sidebar: 
        # Nama pengguna dan log keluar sudah dipaparkan di main() sebelum panggil display_sidebar
        # st.markdown(f"#### 👤 Pengguna: {username}") # Boleh alih keluar jika sudah ada di main
        # if st.button("🚪 Log Keluar", use_container_width=True, key="logout_sidebar_main_app"):
        #     st.session_state.authenticated = False
        #     st.session_state.username = None
        #     st.rerun()
        # st.markdown("---") # Pembahagi selepas log keluar

        st.markdown("## ⚙️ Tetapan & Sesi") 
        st.markdown("---")
        st.markdown("#### Model AI")
        if available_models_list:
            try:
                current_model_index = available_models_list.index(st.session_state.selected_ollama_model)
            except ValueError:
                current_model_index = 0
                if available_models_list: st.session_state.selected_ollama_model = available_models_list[0]
                else: st.session_state.selected_ollama_model = DEFAULT_OLLAMA_MODEL
            selected_model_ui = st.selectbox(
                "Pilih Model:", options=available_models_list, index=current_model_index, key="model_selector_widget",
                label_visibility="collapsed" 
            )
            if selected_model_ui != st.session_state.selected_ollama_model:
                st.session_state.selected_ollama_model = selected_model_ui
        else:
            st.warning("Tiada model AI ditemui.")
            if "selected_ollama_model" not in st.session_state:
                 st.session_state.selected_ollama_model = DEFAULT_OLLAMA_MODEL
        st.markdown("---")
        st.markdown("#### 💬 Sesi Perbualan")
        session_ids = load_all_session_ids(username) # Guna username
        current_session_for_select = st.session_state.session_id
        options = ["➕ Perbualan Baru"] + session_ids
        try:
            current_session_index = options.index(current_session_for_select) if current_session_for_select != "new" else 0
        except ValueError:
            current_session_index = 0
            st.session_state.session_id = "new"
            st.session_state.chat_history = []
            st.session_state.current_filename_prefix = datetime.now().strftime("%Y%m%d_%H%M%S")
        selected_session_id_ui = st.sidebar.selectbox(
            "Pilih atau Mulakan Sesi:", options, index=current_session_index, key="session_selector_widget",
            label_visibility="collapsed"
        )
        st.markdown("---")
        with st.expander("🗑️ Urus Sesi Lanjutan", expanded=False):
            can_delete_current = st.session_state.session_id != "new" and st.session_state.session_id in session_ids
            if can_delete_current:
                if st.button(f"Padam Sesi Semasa: {st.session_state.session_id}", key="delete_current_btn", type="secondary", use_container_width=True):
                    if delete_chat_session_file(username, st.session_state.session_id): # Guna username
                        st.session_state.session_id = "new"; st.session_state.chat_history = []
                        st.session_state.current_filename_prefix = datetime.now().strftime("%Y%m%d_%H%M%S")
                        st.session_state.show_confirm_delete_all_button = False
                        st.session_state.chat_page_num = 1
                        st.rerun()
            if session_ids: # Hanya tunjukkan jika ada sesi untuk pengguna ini
                if not st.session_state.show_confirm_delete_all_button:
                    if st.button("Padam Semua Sesi Saya", key="ask_delete_all_btn", use_container_width=True): # Ubah label
                        st.session_state.show_confirm_delete_all_button = True
                        st.rerun()
                if st.session_state.show_confirm_delete_all_button:
                    st.warning("ANDA PASTI MAHU MEMADAM SEMUA SESI ANDA?") # Ubah label
                    col1, col2 = st.columns(2)
                    with col1:
                        if st.button("YA, PADAM SEMUA MILIK SAYA", key="confirm_delete_all_btn", type="primary", use_container_width=True): # Ubah label
                            if delete_all_chat_sessions(username): # Guna username
                                st.session_state.session_id = "new"; st.session_state.chat_history = []
                                st.session_state.current_filename_prefix = datetime.now().strftime("%Y%m%d_%H%M%S")
                                st.session_state.show_confirm_delete_all_button = False
                                st.session_state.chat_page_num = 1
                                st.rerun()
                    with col2:
                        if st.button("BATAL", key="cancel_delete_all_btn", use_container_width=True):
                            st.session_state.show_confirm_delete_all_button = False
                            st.rerun()
            else:
                st.caption("Tiada sesi untuk dipadam.")
                st.session_state.show_confirm_delete_all_button = False
    return selected_session_id_ui

def handle_session_logic(username, selected_session_id_from_ui): # Tambah username
    if selected_session_id_from_ui == "➕ Perbualan Baru":
        if st.session_state.session_id != "new": 
            st.session_state.session_id = "new"
            st.session_state.chat_history = []
            st.session_state.current_filename_prefix = datetime.now().strftime("%Y%m%d_%H%M%S")
            st.session_state.chat_page_num = 1
    elif st.session_state.session_id != selected_session_id_from_ui: 
        st.session_state.chat_history = load_chat_session(username, selected_session_id_from_ui) # Guna username
        st.session_state.session_id = selected_session_id_from_ui
        st.session_state.current_filename_prefix = selected_session_id_from_ui 
        st.session_state.chat_page_num = 1

def display_chat_messages_paginated():
    # ... (Kod sedia ada anda, tidak berubah) ...
    if not st.session_state.chat_history:
        st.info("💬 Mulakan perbualan dengan menaip di bawah atau muat naik fail untuk analisis.")
        return

    page_size = 10
    total_messages = len(st.session_state.chat_history)
    max_page = (total_messages + page_size - 1) // page_size if total_messages > 0 else 1
    
    if st.session_state.chat_page_num > max_page: st.session_state.chat_page_num = max_page
    if st.session_state.chat_page_num < 1: st.session_state.chat_page_num = 1
    
    reversed_history = st.session_state.chat_history[::-1]
    start_index = (st.session_state.chat_page_num - 1) * page_size
    end_index = start_index + page_size
    messages_to_display = reversed_history[start_index:end_index][::-1] 

    for msg in messages_to_display:
        with st.chat_message(msg["role"]):
            thinking_process_text = msg.get("thinking_process", "").strip()
            if msg["role"] == "assistant" and thinking_process_text:
                with st.expander("Tunjukkan Proses Pemikiran AI", expanded=False):
                    st.markdown(thinking_process_text)
            main_content_text = msg.get("content", "").strip()
            if main_content_text:
                st.markdown(main_content_text)
            elif not thinking_process_text: 
                st.markdown("*(Tiada respons kandungan)*")
            if msg["role"] == "assistant" and "time_taken" in msg and msg["time_taken"] is not None:
                st.caption(f"⏱️ {msg['time_taken']:.2f}s")

    if max_page > 1:
        cols_pagination = st.columns([1, 3, 1]) 
        with cols_pagination[1]:
            page_num_ui = st.slider(
                "Halaman:", 
                min_value=1, 
                max_value=max_page, 
                value=st.session_state.chat_page_num, 
                key="chat_page_slider",
            )
            if page_num_ui != st.session_state.chat_page_num:
                st.session_state.chat_page_num = page_num_ui
                st.rerun() 
    else:
        st.session_state.chat_page_num = 1

def display_export_options():
    # ... (Kod sedia ada anda, tidak berubah) ...
    if not st.session_state.chat_history:
        return
    st.markdown("---") 
    with st.expander("📤 Eksport Perbualan", expanded=False): 
        col_export1, col_export2 = st.columns(2)
        with col_export1:
            export_content_choice = st.radio(
                "Kandungan:",
                ["Pembantu Sahaja", "Pengguna Sahaja", "Keseluruhan Perbualan"],
                index=2, key="export_content_radio"
            )
        with col_export2:
            export_format_choice = st.selectbox("Format:", [
                "Pilih format", "Word (.docx)", "Teks (.txt)", "PDF (.pdf)",
                "Excel (.xlsx)", "PowerPoint (.pptx)"
            ], key="export_format_select")
        custom_filename_prefix_ui = st.text_input(
            "Nama fail awalan:",
            st.session_state.current_filename_prefix,
            key="filename_prefix_input"
        )
        if st.button("📁 Eksport Sekarang", key="export_main_button", type="primary", use_container_width=True):
            if export_format_choice == "Pilih format":
                st.warning("Sila pilih format eksport."); return
            filename_base = custom_filename_prefix_ui 
            include_user = "Pengguna" in export_content_choice or "Keseluruhan" in export_content_choice
            include_assistant = "Pembantu" in export_content_choice or "Keseluruhan" in export_content_choice
            text_for_common_formats = format_conversation_text(st.session_state.chat_history, include_user, include_assistant)
            history_for_excel_pptx = [
                msg for msg in st.session_state.chat_history 
                if (include_user and msg["role"] == "user") or \
                   (include_assistant and msg["role"] == "assistant")
            ]
            if not history_for_excel_pptx and (export_format_choice in ["Excel (.xlsx)", "PowerPoint (.pptx)"]):
                st.warning(f"Tiada mesej '{export_content_choice.lower().replace(' keseluruhan perbualan', '')}' untuk dieksport.")
                return
            success, exported_filename = False, ""
            actions = {
                "Word (.docx)": (save_to_word, text_for_common_formats, f"{filename_base}.docx"),
                "Teks (.txt)": (save_to_txt, text_for_common_formats, f"{filename_base}.txt"),
                "PDF (.pdf)": (save_to_pdf, text_for_common_formats, f"{filename_base}.pdf"),
                "Excel (.xlsx)": (save_to_excel, history_for_excel_pptx, f"{filename_base}.xlsx"), 
                "PowerPoint (.pptx)": (save_to_pptx, history_for_excel_pptx, f"{filename_base}.pptx") 
            }
            if export_format_choice in actions:
                func, data_to_export, fname = actions[export_format_choice]
                if not data_to_export:
                     st.warning(f"Tiada kandungan untuk dieksport ke {export_format_choice}.")
                     return
                success = func(data_to_export, fname)
                exported_filename = fname
            if success and exported_filename:
                st.success(f"Fail disimpan: {exported_filename}")
                try:
                    with open(exported_filename, "rb") as f_download:
                        st.download_button(
                            "📥 Muat Turun", data=f_download, file_name=exported_filename, 
                            key=f"download_btn_{exported_filename.replace('.', '_')}_{time.time()}",
                            use_container_width=True
                        )
                except FileNotFoundError: st.error(f"Gagal mencari {exported_filename} untuk dimuat turun.")
                except Exception as e: st.error(f"Ralat muat turun: {e}")

# --- PERUBAHAN PADA FUNGSI main ---
def main():
    st.set_page_config(page_title="DFK Stembot", layout="wide", initial_sidebar_state="expanded", page_icon="🤖")
    
    # Paparkan logo (menggunakan LOGO_PATH)
    if os.path.exists(LOGO_PATH):
        col_logo_space1, col_logo, col_logo_space2 = st.columns([1.6, 2, 1]) # Laraskan nisbah jika perlu
        with col_logo:
            try:
                st.image(LOGO_PATH, width=250) # Laraskan saiz jika perlu
            except Exception as e:
                st.warning(f"Gagal memaparkan logo: {e}")
    # else: # Tidak perlu warning di sini jika logo adalah pilihan
    #     st.warning(f"Fail logo tidak ditemui di: {LOGO_PATH}") 
    
    # Semak status log masuk
    if not authentication_ui():
        return # Hentikan pelaksanaan jika tidak disahkan
    
    # Jika disahkan, dapatkan username dari session_state
    current_username = st.session_state.username

    # Paparkan nama pengguna dan butang log keluar di sidebar
    with st.sidebar:
        st.markdown(f"#### 👤 Pengguna: {current_username}")
        if st.button("🚪 Log Keluar", use_container_width=True, key="main_logout_button"):
            # Reset state berkaitan sesi pengguna
            st.session_state.authenticated = False
            st.session_state.username = None
            st.session_state.session_id = "new" # Reset ke sesi baru
            st.session_state.chat_history = []
            # Anda mungkin mahu reset state lain yang spesifik pengguna di sini
            st.rerun()
        st.markdown("---") # Pembahagi selepas log keluar

    # Kod utama chatbot selepas ini
    st.markdown("<h1 style='text-align: center; margin-top: 0px; margin-bottom: 10px;'>🤖 DFK Stembot</h1>", unsafe_allow_html=True)
    
    available_ollama_models = get_ollama_models_cached()
    if not available_ollama_models:
        st.error("Tidak dapat memuatkan senarai model dari Ollama.")
    
    initialize_session_state(available_ollama_models) # Inisialisasi state umum

    if st.session_state.selected_ollama_model:
        st.markdown(f"<p style='text-align: center; color: grey; margin-bottom: 20px;'>Model Aktif: <b>{st.session_state.selected_ollama_model.split(':')[0]}</b></p>", unsafe_allow_html=True)
    else:
        st.markdown("<p style='text-align: center; color: red; margin-bottom: 20px;'>Model tidak dipilih</p>", unsafe_allow_html=True)
    
    st.markdown("---")

    selected_session_id_from_ui = display_sidebar(available_ollama_models, current_username) # Hantar username
    handle_session_logic(current_username, selected_session_id_from_ui) # Hantar username
    
    with st.sidebar:
        # st.markdown("---") # Pembahagi sudah ada di atas
        st.markdown("#### 📎 Muat Naik & Analisis Fail")
        uploader_key = f"file_uploader_{st.session_state.uploader_key_counter}"
        uploaded_file = st.file_uploader(
            "Pilih fail (Imej, PDF, DOCX, TXT):", 
            type=['png', 'jpg', 'jpeg', 'gif', 'pdf', 'txt', 'docx'],
            key=uploader_key,
            label_visibility="collapsed"
        )

        if uploaded_file is not None:
            with st.spinner(f"Memproses '{uploaded_file.name}'..."):
                extracted_text = extract_text_from_file(uploaded_file)
            
            if extracted_text:
                st.success(f"Teks diekstrak dari '{uploaded_file.name}'.")
                file_content_message = f"Kandungan dari fail '{uploaded_file.name}':\n\n{extracted_text}"
                st.session_state.chat_history.append({"role": "user", "content": file_content_message})
                
                with st.spinner(f"Menganalisis kandungan fail..."):
                    assistant_response, thinking_text, gen_time = query_ollama_non_stream(
                        file_content_message,
                        st.session_state.chat_history, 
                        st.session_state.selected_ollama_model
                    )
                st.session_state.chat_history.append({
                    "role": "assistant", 
                    "content": assistant_response,
                    "thinking_process": thinking_text,
                    "time_taken": gen_time
                })
                if st.session_state.session_id == "new":
                    st.session_state.session_id = st.session_state.current_filename_prefix
                save_chat_session(current_username, st.session_state.session_id, st.session_state.chat_history) # Guna username
            
            elif extracted_text is None: 
                pass 
            else: 
                st.warning(f"Tiada teks diekstrak dari '{uploaded_file.name}'.")
            st.session_state.uploader_key_counter += 1
            st.rerun()

    chat_container = st.container() 
    with chat_container:
        display_chat_messages_paginated()

    user_input = st.chat_input(f"Tanya {st.session_state.selected_ollama_model.split(':')[0].capitalize()}...")

    if user_input:
        st.session_state.chat_history.append({"role": "user", "content": user_input})
        with st.spinner(f"Sedang berfikir..."): 
            assistant_response_text, thinking_text, generation_time = query_ollama_non_stream(
                user_input, 
                st.session_state.chat_history, 
                st.session_state.selected_ollama_model
            )
        st.session_state.chat_history.append({
            "role": "assistant", 
            "content": assistant_response_text,
            "thinking_process": thinking_text,
            "time_taken": generation_time
        })
        if st.session_state.session_id == "new":
            st.session_state.session_id = st.session_state.current_filename_prefix
        save_chat_session(current_username, st.session_state.session_id, st.session_state.chat_history) # Guna username
        
        total_messages = len(st.session_state.chat_history)
        page_size = 10 
        st.session_state.chat_page_num = (total_messages + page_size - 1) // page_size if total_messages > 0 else 1
        st.rerun()

    display_export_options()

if __name__ == "__main__":
    main()
        # Cadangan untuk fail requirements.txt:
    #
    # streamlit>=1.17.0
    # requests
    # python-docx
    # fpdf2
    # pandas
    # openpyxl
    # python-pptx
    # Pillow
    # pytesseract
    # PyMuPDF
    # bcrypt
    # Pasang menggunakan: pip install -r requirements.txt
    # Juga, pastikan Tesseract OCR dipasang pada sistem anda.
